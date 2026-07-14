#!/usr/bin/env python3
"""
PrimeNet — Engineering Deep-Dive deck generator.

Design goals:
- Audience: engineers. In-depth architecture / data-flow / algorithm content.
- Built as a *demo companion*: the presenter constantly switches between the
  live PrimeNet tool and these slides, so the visual language mirrors the app
  (constellation dark theme) and every major topic carries an explicit
  "SWITCH TO TOOL" cue with an exact click-path, plus "BACK TO DECK" recaps.
"""

import math
import random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette ----
BG        = RGBColor(0x0F, 0x17, 0x22)   # constellation background
BG2       = RGBColor(0x0B, 0x11, 0x1A)   # deeper panel for dividers
PANEL     = RGBColor(0x18, 0x22, 0x30)
PANEL2    = RGBColor(0x1B, 0x27, 0x36)
PANEL3    = RGBColor(0x22, 0x32, 0x46)
BORDER    = RGBColor(0x30, 0x42, 0x58)
TEXT      = RGBColor(0xE8, 0xEE, 0xF7)
MUTED     = RGBColor(0xA9, 0xB7, 0xC9)
FAINT     = RGBColor(0x74, 0x86, 0x9C)
ACCENT    = RGBColor(0x8B, 0xC1, 0xFF)   # link blue
TEAL      = RGBColor(0x3D, 0xD6, 0xB0)
GOLD      = RGBColor(0xF2, 0xB8, 0x4B)   # DEMO / switch cue
PURPLE    = RGBColor(0xB0, 0x8C, 0xFF)
CRIT      = RGBColor(0xE5, 0x54, 0x8A)   # pinkish-red for accents
ORANGE    = RGBColor(0xE9, 0x91, 0x3A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]

FONT   = "Segoe UI"
FONT_L = "Segoe UI Light"
MONO   = "Consolas"

_slide_no = 0


# --------------------------------------------------------------- helpers -----
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_shadow(shape):
    # python-pptx auto-shadows some autoshapes; strip it.
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    for tag in ('a:effectLst', 'a:effectDag'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    spPr.append(spPr.makeelement(qn('a:effectLst'), {}))


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.06,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _no_shadow(sp)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def line(slide, x1, y1, x2, y2, color=BORDER, w=1.0, dash=None):
    ln = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    if dash:
        d = ln.line._get_or_add_ln()
        pd = d.makeelement(qn('a:prstDash'), {'val': dash})
        d.append(pd)
    _no_shadow(ln)
    return ln


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.06, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font, italic)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for m in ('left', 'right', 'top', 'bottom'):
        setattr(tf, f'margin_{m}', 0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, *rest) in para:
            fname = rest[0] if len(rest) > 0 and rest[0] else FONT
            italic = rest[1] if len(rest) > 1 else False
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = fname
            r.font.italic = italic
    return tb


def slide_base(divider=False):
    s = prs.slides.add_slide(BLANK)
    rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, fill=(BG2 if divider else BG),
         shape=MSO_SHAPE.RECTANGLE)
    return s


def constellation(slide, n=46, seed=1, box=(0, 0, SW, SH), link=1.7, alpha_dim=True):
    """Sparse star field with a few connecting lines — the PrimeNet motif."""
    random.seed(seed)
    x0, y0, x1, y1 = box
    pts = []
    for _ in range(n):
        px = random.uniform(x0, x1)
        py = random.uniform(y0, y1)
        pts.append((px, py))
    # links between near points
    for i, a in enumerate(pts):
        for b in pts[i + 1:]:
            d = math.hypot(a[0] - b[0], a[1] - b[1])
            if d < link and random.random() < 0.5:
                line(slide, a[0], a[1], b[0], b[1], color=RGBColor(0x22, 0x30, 0x44), w=0.75)
    for (px, py) in pts:
        r = random.choice([0.02, 0.03, 0.045, 0.03])
        c = random.choice([ACCENT, TEAL, RGBColor(0x3A, 0x4A, 0x60), MUTED, ACCENT])
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px), Inches(py),
                                     Inches(r), Inches(r))
        _set_fill(dot, c)
        _no_shadow(dot)


def footer(slide, section=""):
    global _slide_no
    _slide_no += 1
    line(slide, 0.55, SH - 0.52, SW - 0.55, SH - 0.52, color=RGBColor(0x22, 0x30, 0x42), w=0.75)
    text(slide, 0.55, SH - 0.47, 6.5, 0.3,
         [[("PrimeNet", 9, ACCENT, True), ("  ·  Engineering Deep-Dive", 9, FAINT, False)]],
         anchor=MSO_ANCHOR.MIDDLE)
    if section:
        text(slide, SW / 2 - 2.5, SH - 0.47, 5.0, 0.3, [[(section, 9, FAINT, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, SW - 2.05, SH - 0.47, 1.5, 0.3, [[(f"{_slide_no:02d}", 9, MUTED, True)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(slide, kicker, title, kicker_color=ACCENT):
    rect(slide, 0.55, 0.5, 0.09, 0.62, fill=kicker_color, shape=MSO_SHAPE.RECTANGLE)
    text(slide, 0.78, 0.46, 11.8, 0.3,
         [[(kicker.upper(), 11, kicker_color, True)]])
    text(slide, 0.77, 0.7, 12.0, 0.6, [[(title, 25, TEXT, True, FONT_L)]])


def chip(slide, x, y, w, label, color, txtcolor=None, h=0.34, size=10.5, bold=True):
    c = rect(slide, x, y, w, h, fill=None, line=color, line_w=1.25, radius=0.5)
    text(slide, x, y - 0.005, w, h, [[(label, size, txtcolor or color, bold)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def bullets(slide, x, y, w, h, items, size=13, gap=7, lead=ACCENT, tcolor=TEXT,
            line_spacing=1.05):
    """items: list of (bold_lead, rest) or plain string."""
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead_t, rest = it
            paras.append([("▸  ", size, lead, True),
                          (lead_t, size, WHITE, True),
                          (rest, size, tcolor, False)])
        else:
            paras.append([("▸  ", size, lead, True), (it, size, tcolor, False)])
    text(slide, x, y, w, h, paras, space_after=gap, line_spacing=line_spacing)


def demo_cue(slide, y, title, steps, watch=None):
    """Amber 'switch to the live tool' band — the core companion-deck device."""
    x, w = 0.55, SW - 1.10
    h = 1.3 if watch else 1.06
    rect(slide, x, y, w, h, fill=RGBColor(0x22, 0x1C, 0x0E), line=GOLD, line_w=1.25, radius=0.05)
    rect(slide, x, y, 0.09, h, fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
    text(slide, x + 0.28, y + 0.12, 6.5, 0.3,
         [[("⟳  SWITCH TO TOOL", 12, GOLD, True), ("   — live demo checkpoint", 10.5, RGBColor(0xC9,0xB0,0x82), False, FONT, True)]])
    text(slide, x + 0.28, y + 0.44, w - 0.6, 0.3, [[(title, 12, TEXT, True)]])
    text(slide, x + 0.28, y + 0.71, w - 0.6, 0.4,
         [[("Click path:  ", 10.5, GOLD, True), (steps, 10.5, MUTED, False, MONO)]])
    if watch:
        text(slide, x + 0.28, y + 0.98, w - 0.6, 0.3,
             [[("Point out:  ", 10.5, TEAL, True), (watch, 10.5, MUTED, False)]])
    return h


# ============================================================ SLIDE 1: TITLE
s = slide_base(divider=True)
constellation(s, n=70, seed=7, link=1.9)
# center glow panel
text(s, 0.9, 1.5, 11.5, 0.4, [[("ENGINEERING DEEP-DIVE", 15, ACCENT, True)]])
text(s, 0.85, 1.95, 11.6, 1.5, [[("PrimeNet", 62, WHITE, True, FONT_L)]])
text(s, 0.9, 3.15, 11.5, 0.6,
     [[("Radio network performance & configuration platform — under the hood", 18, MUTED, False, FONT_L)]])
# meta chips
chips = [("Flask modular monolith", ACCENT), ("Nokia + Huawei", TEAL),
         ("2G · 3G · 4G · 5G", PURPLE), ("37 feature modules", GOLD)]
cx = 0.9
for lab, col in chips:
    wch = 0.24 + 0.105 * len(lab)
    chip(s, cx, 4.05, wch, lab, col, h=0.42, size=11.5)
    cx += wch + 0.22
# how-to-read note
rect(s, 0.9, 5.05, 11.5, 1.05, fill=RGBColor(0x14,0x1E,0x2B), line=GOLD, line_w=1.0, radius=0.05)
rect(s, 0.9, 5.05, 0.09, 1.05, fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
text(s, 1.15, 5.2, 11.0, 0.8,
     [[("⟳  This is a companion deck.  ", 12.5, GOLD, True),
       ("Amber bands mark ", 12, MUTED, False),
       ("live-tool checkpoints", 12, TEXT, True),
       (" — keep PrimeNet open in a second window and switch to it whenever you see one.", 12, MUTED, False)],
      [("The slides carry the engineering that the UI doesn't show; the tool carries the proof.", 11, FAINT, False, FONT, True)]],
     space_after=5)
text(s, 0.9, 6.55, 11.5, 0.3, [[("Presenter: Malek Mohammad   ·   Audience: RAN / platform engineers", 11, FAINT, False)]])

# ============================================================ SLIDE 2: HOW TO USE
s = slide_base()
header(s, "Presenter guide", "How to run this session")
# two-window setup illustration
lx = 0.55
text(s, lx, 1.35, 7.2, 0.4, [[("The two-window rhythm", 14, TEAL, True)]])
bullets(s, lx, 1.75, 7.2, 3.0, [
    ("Deck = the “why”. ", "Architecture, data flow, algorithms, and design trade-offs that never appear on screen."),
    ("Tool = the “proof”. ", "Every claim is demonstrable live; the deck tells you exactly where to click."),
    ("Ideal layout: ", "two monitors, or deck on the left half / PrimeNet on the right half."),
    ("Cadence: ", "explain a subsystem → hit its amber checkpoint → return on the recap line."),
    ("Login once up front ", "so activation and session are warm before the first demo."),
], size=13, gap=9)
# legend card
rx, rw = 8.15, 4.65
rect(s, rx, 1.35, rw, 4.55, fill=PANEL, line=BORDER, radius=0.04)
text(s, rx + 0.3, 1.55, rw - 0.6, 0.3, [[("SLIDE LEGEND", 11, MUTED, True)]])
legend = [
    (GOLD,  "⟳ SWITCH TO TOOL", "Stop talking, start clicking. Click-path is monospaced."),
    (TEAL,  "Point out",        "The one thing the audience must notice on screen."),
    (ACCENT,"▸ blue bullets",   "Load-bearing engineering detail."),
    (PURPLE,"◆ diagrams",       "Data-flow / lifecycle maps to anchor the mental model."),
]
yy = 2.0
for col, lab, desc in legend:
    rect(s, rx + 0.3, yy + 0.05, 0.16, 0.16, fill=col, shape=MSO_SHAPE.OVAL)
    text(s, rx + 0.6, yy - 0.05, rw - 0.9, 0.3, [[(lab, 12, col, True)]])
    text(s, rx + 0.6, yy + 0.26, rw - 0.9, 0.5, [[(desc, 10.5, MUTED, False)]])
    yy += 0.86
text(s, rx + 0.3, yy + 0.05, rw - 0.6, 0.6,
     [[("If a demo env is offline, the deck still stands alone — just narrate the checkpoint.", 10.5, FAINT, False, FONT, True)]])
footer(s, "Presenter guide")

# ============================================================ SLIDE 3: AGENDA
s = slide_base()
header(s, "Roadmap", "What we'll cover")
agenda = [
    ("01", "System at a glance", "Scale, shape, and the modular-monolith bet", ACCENT),
    ("02", "Architecture & a module's anatomy", "Blueprints, folders, shared core", TEAL),
    ("03", "Request lifecycle", "The before_request middleware chain & security", PURPLE),
    ("04", "Data platform & ETL", "SFTP pull → raw taxonomy → SQLite load", GOLD),
    ("05", "Orchestration & scheduling", "Daily/hourly/watcher, APScheduler", ORANGE),
    ("06", "CM Extractor & write-back", "Nokia CM Open API, Huawei MML, RET", ACCENT),
    ("07", "Radio analytics engine", "The shared scoring model & insight catalog", TEAL),
    ("08", "Frontend & deployment", "Theme system, Docker, extending PrimeNet", PURPLE),
]
colw, rowh = 5.9, 1.18
x0, y0 = 0.6, 1.45
for i, (n, t, d, col) in enumerate(agenda):
    cx = x0 + (i % 2) * (colw + 0.35)
    cy = y0 + (i // 2) * (rowh + 0.12)
    rect(s, cx, cy, colw, rowh, fill=PANEL, line=BORDER, radius=0.06)
    rect(s, cx, cy, 0.07, rowh, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx + 0.28, cy + 0.12, 1.2, 0.9, [[(n, 30, col, True, FONT_L)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 1.35, cy + 0.19, colw - 1.5, 0.4, [[(t, 14.5, TEXT, True)]])
    text(s, cx + 1.35, cy + 0.62, colw - 1.5, 0.4, [[(d, 11, MUTED, False)]])
footer(s, "Roadmap")

# ============================================================ SLIDE 4: AT A GLANCE
s = slide_base()
header(s, "Section 01", "System at a glance")
stats = [
    ("37", "feature modules", "each a self-contained Flask blueprint", ACCENT),
    ("2", "RAN vendors", "Nokia & Huawei, one shared UI", TEAL),
    ("4", "generations", "2G / 3G / 4G / 5G KPIs & CM", PURPLE),
    ("3", "SFTP sources", "Nokia PM · Huawei PM · Metadata", GOLD),
]
cw = 2.92
for i, (big, lab, sub, col) in enumerate(stats):
    cx = 0.6 + i * (cw + 0.13)
    rect(s, cx, 1.4, cw, 1.7, fill=PANEL, line=BORDER, radius=0.06)
    text(s, cx + 0.25, 1.5, cw - 0.5, 0.9, [[(big, 44, col, True, FONT_L)]])
    text(s, cx + 0.27, 2.42, cw - 0.5, 0.3, [[(lab, 13, TEXT, True)]])
    text(s, cx + 0.27, 2.72, cw - 0.5, 0.4, [[(sub, 10, MUTED, False)]])
# narrative
text(s, 0.6, 3.35, 12.1, 0.4, [[("The shape of the system", 14, TEAL, True)]])
bullets(s, 0.6, 3.72, 6.05, 3.0, [
    ("One process, many modules. ", "A modular monolith: 37 blueprints registered in app.py, one SQLite-backed runtime, one auth/session layer."),
    ("Vendor + RAT are data, not forks. ", "Nokia/Huawei and 2G–5G differences live in adapters and a path taxonomy, not parallel apps."),
    ("Read-heavy analytics, guarded writes. ", "PM analytics are read-only; CM write-back (RET, mass-modify) is explicit and audited."),
], size=12.5, gap=8)
bullets(s, 6.85, 3.72, 5.85, 3.0, [
    ("Why a monolith? ", "Small ops team, shared SQLite data, no per-module network hops — deploy is a single container."),
    ("Cohesion via convention. ", "Every module follows the same folder + blueprint contract, so 37 modules stay legible."),
    ("Escape hatches exist. ", "Heavy CM/FM work delegates to core/ clients; the pipeline runs as separate orchestrated jobs."),
], size=12.5, gap=8, lead=TEAL)
footer(s, "01 · Overview")

# ============================================================ SLIDE 5: ARCHITECTURE
s = slide_base()
header(s, "Section 02", "Architecture: the modular monolith")
# layered diagram
def band(x, y, w, h, title, items, col, tsize=11.5):
    rect(s, x, y, w, h, fill=PANEL, line=col, line_w=1.25, radius=0.05)
    rect(s, x, y, 0.07, h, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.22, y + 0.09, w - 0.35, 0.3, [[(title, tsize, col, True)]])
    text(s, x + 0.22, y + 0.42, w - 0.35, h - 0.5, [[(items, 10, MUTED, False)]], line_spacing=1.02)

band(1.6, 1.35, 10.1, 0.86, "App shell — app.py",
     "Loads .env → activation gate → registers 37 blueprints → before/after_request hooks → error handlers → health probe", ACCENT)
# arrow
line(s, 6.65, 2.21, 6.65, 2.45, color=BORDER, w=1.5)
band(0.6, 2.5, 5.75, 1.55, "Feature modules  ·  modules/<name>/",
     "routes.py  (blueprint)\ntemplates/   (Jinja pages)\nstatic/      (module CSS/JS)\nlogic.py     (module compute)\n\n37 modules — Performance, CM Extractor,\nRET, Sector Health, SON, Fault Mgmt …", TEAL, tsize=12)
band(6.55, 2.5, 6.15, 1.55, "Shared infrastructure",
     "routes/auth_routes.py      login / session\ncore/                       vendor clients, scoring, radio,\n                            cm_extractor, activation, licensing\nutils/  ·  db/runtime.py    SQLite access, path constants\ncore/module_access.py       RBAC visibility rules", PURPLE, tsize=12)
line(s, 6.65, 4.05, 6.65, 4.3, color=BORDER, w=1.5)
band(0.6, 4.35, 5.75, 1.35, "Data platform",
     "SQLite files under databases/\n   cells / groups / metadata / admin\nCanonical taxonomy: pipeline/paths.py\nActivation-gated connections (db/runtime.py)", GOLD, tsize=12)
band(6.55, 4.35, 6.15, 1.35, "ETL pipeline  ·  pipeline/",
     "pull/  →  raw/{vendor}/{domain}/{rat}/{tf}\nload/  →  databases/…\norchestrators/ (daily · hourly · watcher)\nAPScheduler drives recurring jobs", ORANGE, tsize=12)
text(s, 0.6, 5.9, 12.1, 0.6,
     [[("Contract, not framework:  ", 12, TEAL, True),
       ("a new capability is a folder that follows the module shape and one register_blueprint() line. That single convention is what keeps 37 modules maintainable by a small team.", 12, MUTED, False)]])
footer(s, "02 · Architecture")

# ============================================================ SLIDE 6: MODULE ANATOMY
s = slide_base()
header(s, "Section 02", "Anatomy of a module")
# left: tree
rect(s, 0.55, 1.4, 5.7, 3.9, fill=RGBColor(0x0C,0x13,0x1D), line=BORDER, radius=0.04)
tree = [
    ("modules/capacity_hotspots/", ACCENT, True),
    ("├─ __init__.py", MUTED, False),
    ("├─ routes.py", TEAL, False),
    ("│    @login_required", FAINT, False),
    ("│    capacity_hotspots_bp = Blueprint(…)", FAINT, False),
    ("├─ logic.py            # compute / scoring", TEAL, False),
    ("├─ templates/", MUTED, False),
    ("│    └─ capacity_hotspots.html", FAINT, False),
    ("└─ static/", MUTED, False),
    ("     ├─ capacity_hotspots.css", FAINT, False),
    ("     └─ capacity_hotspots.js", FAINT, False),
]
yy = 1.62
for t, c, b in tree:
    text(s, 0.8, yy, 5.3, 0.3, [[(t, 12, c, b, MONO)]])
    yy += 0.335
# right: the contract
text(s, 6.55, 1.4, 6.1, 0.3, [[("The module contract", 14, TEAL, True)]])
bullets(s, 6.55, 1.78, 6.15, 3.0, [
    ("Blueprint per module. ", "Own url_prefix, template_folder=\"templates\", module-local static_folder."),
    ("Auth by decorator. ", "Copy the login_required pattern; session rides a session_token cookie."),
    ("Thin routes. ", "routes.py validates + delegates; real work sits in logic.py or core/."),
    ("Visibility is declared. ", "core/module_access.py maps each route to all / admin / admin_or_noc."),
    ("Registered once. ", "Import + app.register_blueprint(...) in app.py — the only global touch-point."),
], size=12, gap=7)
footer(s, "02 · Architecture")
demo_cue(s, 5.5,
         "Open any module and show the identical shell — header, filters, dark-mode toggle.",
         "Dashboard → Radio Optimization → Capacity Hotspots",
         watch="Same chrome across modules = the contract paying off. Toggle dark mode; it persists.")

# ============================================================ SLIDE 7: REQUEST LIFECYCLE
s = slide_base()
header(s, "Section 03", "Request lifecycle: the middleware chain")
text(s, 0.55, 1.28, 12.2, 0.3,
     [[("Every request threads four ", 12, MUTED, False), ("@app.before_request", 12, TEAL, True, MONO),
       (" guards before a blueprint ever runs:", 12, MUTED, False)]])
steps = [
    ("1", "Activation gate", "enforce_monthly_operator_activation — locked → /activation (or 403 on /api).", ACCENT),
    ("2", "Input safety", "Validate + sanitize args / form / JSON; reject malformed or oversized bodies (413/400).", TEAL),
    ("3", "CSRF origin", "State-changing + cookie auth ⇒ Origin/Referer must be same-origin, else 403.", PURPLE),
    ("4", "Password rotation", "Force expired credentials to /profile before any other page.", GOLD),
]
x = 0.55
bw = 3.0
for i, (n, t, d, col) in enumerate(steps):
    cx = x + i * (bw + 0.1)
    rect(s, cx, 1.72, bw, 1.95, fill=PANEL, line=col, line_w=1.25, radius=0.05)
    text(s, cx + 0.22, 1.85, 0.8, 0.6, [[(n, 26, col, True, FONT_L)]])
    text(s, cx + 0.22, 2.5, bw - 0.4, 0.35, [[(t, 12.5, TEXT, True)]])
    text(s, cx + 0.22, 2.85, bw - 0.4, 0.9, [[(d, 10, MUTED, False)]])
    if i < 3:
        text(s, cx + bw - 0.03, 2.4, 0.2, 0.4, [[("→", 16, FAINT, True)]])
# after_request
rect(s, 0.55, 3.95, 12.2, 0.95, fill=PANEL2, line=BORDER, radius=0.05)
text(s, 0.8, 4.05, 4.0, 0.3, [[("@app.after_request", 12.5, ORANGE, True, MONO)]])
text(s, 0.8, 4.38, 11.9, 0.5,
     [[("Security headers on every response:  ", 11, MUTED, False),
       ("CSP", 11, TEXT, True), (" (self + unpkg + OSM/ArcGIS tiles), ", 10.5, MUTED, False),
       ("X-Content-Type-Options", 11, TEXT, True), (", ", 10.5, MUTED, False),
       ("X-Frame-Options", 11, TEXT, True), (", ", 10.5, MUTED, False),
       ("Referrer-Policy", 11, TEXT, True), (", ", 10.5, MUTED, False),
       ("Permissions-Policy", 11, TEXT, True), (", optional HSTS.", 10.5, MUTED, False)]])
text(s, 0.55, 5.05, 12.2, 0.7,
     [[("Engineering note:  ", 11.5, TEAL, True),
       ("access logs use a ConciseRequestHandler that strips query strings — KPI selections can be thousands of chars, so the raw URL never hits the log. Fail-closed by default: unknown state ⇒ redirect/deny, not allow.", 11.5, MUTED, False)]])
footer(s, "03 · Request lifecycle")

# ============================================================ SLIDE 8: SECURITY / ACTIVATION
s = slide_base()
header(s, "Section 03", "Security, sessions & activation")
col2 = [
    ("Auth & sessions", TEAL, [
        ("session_token cookie ", "server-side sessions; lifetime via SESSION_LIFETIME_HOURS (default 2h)."),
        ("RBAC ", "three tiers — all / admin (owner) / admin_or_noc — enforced in nav + routes."),
        ("Password policy ", "forced rotation short-circuits the whole app until changed."),
    ]),
    ("Operator activation", GOLD, [
        ("Monthly gate ", "install_sqlite_gate() + activation_gate wrap DB access; unactivated ⇒ 503/redirect."),
        ("License client ", "core/license_client + license_tokens verify the operator token."),
        ("Health-aware ", "/health returns locked / degraded / ok for orchestrators & LBs."),
    ]),
    ("Input & transport", PURPLE, [
        ("Global sanitizer ", "depth/'size/length caps on args, form, JSON via utils/input_safety."),
        ("CSRF ", "same-origin check for cookie-auth mutations."),
        ("Secrets ", "SFTP/CM creds only in local .env — never committed; server-side only."),
    ]),
]
cw = 3.98
for i, (t, col, items) in enumerate(col2):
    cx = 0.55 + i * (cw + 0.13)
    rect(s, cx, 1.4, cw, 4.35, fill=PANEL, line=BORDER, radius=0.04)
    rect(s, cx, 1.4, cw, 0.5, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx + 0.22, 1.4, cw - 0.4, 0.5, [[(t, 13, BG, True)]], anchor=MSO_ANCHOR.MIDDLE)
    yy = 2.05
    for lead, rest in items:
        text(s, cx + 0.22, yy, cw - 0.42, 1.2,
             [[("● ", 11, col, True), (lead, 11.5, WHITE, True)],
              [(rest, 10.5, MUTED, False)]], space_after=2, line_spacing=1.03)
        yy += 1.18
footer(s, "03 · Security")

# ============================================================ SLIDE 9: DATA PLATFORM
s = slide_base()
header(s, "Section 04", "Data platform: SQLite + a path taxonomy")
text(s, 0.55, 1.3, 7.0, 0.4, [[("Why SQLite (on purpose)", 14, TEAL, True)]])
bullets(s, 0.55, 1.7, 6.1, 3.2, [
    ("Zero-ops, file-per-scope. ", "PM KPIs, metadata, users, admin each live in their own .db under databases/."),
    ("Deterministic paths. ", "pipeline/paths.py is the single source of truth for where data lives."),
    ("Portable. ", "NCM_DATA_ROOT relocates the whole data tree to a mounted volume in Docker."),
    ("Gated access. ", "db/runtime.py installs the activation gate before any connection opens."),
    ("Automatic PM columns. ", "KPI column detection is runtime-automatic — no per-vendor mapping to maintain."),
], size=12, gap=7)
# taxonomy card
rx, rw = 6.9, 5.85
rect(s, rx, 1.3, rw, 4.3, fill=RGBColor(0x0C,0x13,0x1D), line=BORDER, radius=0.04)
text(s, rx + 0.28, 1.5, rw - 0.5, 0.3, [[("CANONICAL TAXONOMY", 11, MUTED, True)]])
tax = [
    ("raw/", ACCENT),
    ("  {vendor}/{domain}/{rat}/{timeframe}/", FAINT),
    ("  nokia/cells/4g/hourly/  ← one RAT per folder", TEAL),
    ("  huawei/cells/all/daily/ ← staging, then RAT split", TEAL),
    ("", MUTED),
    ("databases/", ACCENT),
    ("  {domain}/{vendor}/{tech}/{timeframe}/*.db", FAINT),
    ("  cells/nokia/all/hourly/…", TEAL),
    ("", MUTED),
    ("vendors   = nokia · huawei · metadata", MUTED),
    ("domains   = cells · groups · neighbors · …", MUTED),
    ("rats      = 2g · 3g · 4g · 5g", MUTED),
    ("timeframes= hourly · daily · snapshot", MUTED),
]
yy = 1.9
for t, c in tax:
    text(s, rx + 0.3, yy, rw - 0.5, 0.3, [[(t, 11.5, c, (c in (ACCENT,)), MONO)]])
    yy += 0.275
footer(s, "04 · Data platform")

# ============================================================ SLIDE 10: ETL PIPELINE
s = slide_base()
header(s, "Section 04", "ETL pipeline: source → raw → SQLite")
# pipeline flow
stages = [
    ("SOURCES", "3× SFTP servers\nNokia PM · Huawei PM\n· Metadata", ACCENT),
    ("PULL", "pipeline/pull/…\nparamiko SFTP,\nper vendor/rat/tf", TEAL),
    ("RAW", "raw/ taxonomy\none RAT per folder\n(Huawei staged→split)", PURPLE),
    ("LOAD", "pipeline/load/…\nauto column detect,\nupsert into SQLite", GOLD),
    ("DATABASES", "databases/*.db\nread by 37 modules\n+ radio engine", ORANGE),
]
bw = 2.28
for i, (t, d, col) in enumerate(stages):
    cx = 0.55 + i * (bw + 0.18)
    rect(s, cx, 1.5, bw, 1.7, fill=PANEL, line=col, line_w=1.25, radius=0.06)
    rect(s, cx, 1.5, bw, 0.42, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx, 1.5, bw, 0.42, [[(t, 12, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 0.15, 2.0, bw - 0.3, 1.1, [[(d, 10.5, MUTED, False, MONO)]], line_spacing=1.05)
    if i < 4:
        text(s, cx + bw - 0.02, 2.15, 0.22, 0.4, [[("→", 18, FAINT, True)]])
bullets(s, 0.55, 3.5, 12.1, 2.0, [
    ("Idempotent by design. ", "Pull writes into a deterministic tree; load re-runs safely — reprocessing a day just overwrites its rows."),
    ("Vendor asymmetry handled early. ", "Huawei daily exports stage in raw/huawei/{cells,groups}/all/daily, then split by RAT before load."),
    ("Verified, not assumed. ", "core/pipeline_ingest_verify + pm_health check row counts / freshness before data is trusted downstream."),
    ("Prefer orchestrators. ", "pipeline/orchestrators/ over ad-hoc scripts — the taxonomy + entry points are the contract."),
], size=12, gap=7)
footer(s, "04 · ETL")
demo_cue(s, 5.62,
         "Show a live sync run and the resulting freshness / row counts.",
         "Dashboard → Sync   (watch a pull→load cycle, then Network Health for freshness)")

# ============================================================ SLIDE 11: ORCHESTRATION
s = slide_base()
header(s, "Section 05", "Orchestration & scheduling")
cards = [
    ("orchestrate_daily_full", ACCENT, "Full daily cycle: pull every source's daily exports, split RATs, load all domains, verify. The heavy nightly run."),
    ("orchestrate_hourly_full", TEAL, "Lighter hourly cadence for near-real-time KPIs — same pull→load→verify shape, hourly timeframe."),
    ("orchestrate_watcher_cycle", PURPLE, "Watches for freshly-dropped files and ingests on arrival — bridges the gap between scheduled runs."),
]
for i, (t, col, d) in enumerate(cards):
    cx = 0.55 + i * (4.05)
    rect(s, cx, 1.45, 3.9, 1.9, fill=PANEL, line=col, line_w=1.25, radius=0.05)
    rect(s, cx, 1.45, 0.07, 1.9, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx + 0.24, 1.6, 3.6, 0.5, [[(t, 12.5, col, True, MONO)]])
    text(s, cx + 0.24, 2.15, 3.55, 1.1, [[(d, 11, MUTED, False)]], line_spacing=1.06)
bullets(s, 0.55, 3.7, 12.1, 2.0, [
    ("APScheduler in-process. ", "Recurring jobs run inside the app (requirements: APScheduler) — no external cron dependency for the core cadence."),
    ("Config Task Scheduler module. ", "User-facing scheduling of CM jobs (exports, audits) with its own persistence and status UI."),
    ("Windows dev affordance. ", "A live-sync logger terminal + auto-browser open are dev-only, disabled by NCM_CONTAINER in Docker."),
    ("Observability. ", "sync_log entries stream live; /health surfaces DB reachability for external schedulers."),
], size=12, gap=7)
footer(s, "05 · Orchestration")

# ============================================================ SLIDE 12: CM EXTRACTOR
s = slide_base()
header(s, "Section 06", "CM Extractor: reading two vendors")
# Nokia
rect(s, 0.55, 1.4, 6.0, 3.05, fill=PANEL, line=ACCENT, line_w=1.25, radius=0.05)
text(s, 0.8, 1.55, 5.5, 0.3, [[("NOKIA — CM Open API (NetAct)", 13, ACCENT, True)]])
text(s, 0.8, 1.95, 5.5, 2.4, [[
    ("Two REST interfaces, HTTP Basic:", 11, TEXT, True)]], space_after=4)
nk = [
    ("persistency/v1", "read CM — queries, MO lists, parameter reads"),
    ("operations/v1", "Configurator ops — Provision, Export, Import_Export"),
    ("confId", "1 = live · 5 = reference · other = plan (writes)"),
    ("DN paths", "PLMN-…/RNC-42/WBTS-176/WCEL-1  ·  class NOKRNC:WCEL"),
    ("Throttled", "batch size + delay + retries to dodge NetAct 429s"),
]
yy = 2.32
for a, b in nk:
    text(s, 0.8, yy, 5.5, 0.4, [[("• ", 11, ACCENT, True), (a, 11, TEXT, True, MONO), ("  " + b, 10, MUTED, False)]])
    yy += 0.4
# Huawei
rect(s, 6.72, 1.4, 6.0, 3.05, fill=PANEL, line=CRIT, line_w=1.25, radius=0.05)
text(s, 6.97, 1.55, 5.5, 0.3, [[("HUAWEI — MML / U2020", 13, CRIT, True)]])
hw = [
    ("MML discovery", "enumerate NEs & MO scope via MML commands"),
    ("huawei_client", "session, command exec, response parsing"),
    ("mml_parser", "structured tables out of raw MML text"),
    ("param dict", "19k scraped reference pages resolve semantics"),
    ("RET / tilt", "RET MOD command formatting + tilt handling"),
]
yy = 2.05
for a, b in hw:
    text(s, 6.97, yy, 5.5, 0.4, [[("• ", 11, CRIT, True), (a, 11, TEXT, True, MONO), ("  " + b, 10, MUTED, False)]])
    yy += 0.45
text(s, 0.55, 4.62, 12.2, 0.5,
     [[("One UI, two protocols.  ", 12, TEAL, True),
       ("core/cm_extractor/ hides the Nokia-REST vs Huawei-MML split behind a common extract → normalize → Excel path; the module layer never sees the difference.", 12, MUTED, False)]])
footer(s, "06 · CM Extractor")
demo_cue(s, 5.55,
         "Run an extraction and open the generated Excel — show the MO/parameter structure.",
         "Dashboard → Configuration → Configuration Data Extractor → pick scope → Extract")

# ============================================================ SLIDE 13: CM WRITE-BACK / RET
s = slide_base()
header(s, "Section 06", "Write-back: RET, mass-modify & audit", kicker_color=GOLD)
bullets(s, 0.55, 1.4, 6.05, 3.6, [
    ("Reads are cheap; writes are governed. ", "Analytics never mutate the network — CM changes go through explicit, reviewable flows."),
    ("RET Management. ", "Read + write electrical tilt across vendors; Huawei RET MOD command formatting and tilt handling are handled centrally."),
    ("Nokia mass-modify. ", "nokia_mass_modify + Excel re-import: export MOs → edit in Excel → validate → write back to a plan confId."),
    ("CM Parameter Audit. ", "Live scanner that diffs current parameters against expected, exportable to Excel."),
    ("Config History. ", "Every applied change is recorded — what/when/who for rollback and forensics."),
], size=12, gap=8)
# write flow
rx, rw = 6.9, 5.85
rect(s, rx, 1.4, rw, 3.55, fill=RGBColor(0x0C,0x13,0x1D), line=GOLD, line_w=1.0, radius=0.04)
text(s, rx + 0.28, 1.55, rw - 0.5, 0.3, [[("SAFE WRITE FLOW", 11, GOLD, True)]])
flow = [
    ("1  Export", "read live MOs (confId 1) → Excel", ACCENT),
    ("2  Edit", "operator edits values in Excel", TEAL),
    ("3  Validate", "re-import checks types / ranges / DNs", PURPLE),
    ("4  Stage", "write to a plan confId (not live)", GOLD),
    ("5  Provision", "operations/v1 applies the plan", ORANGE),
    ("6  Record", "Config History logs the change", CRIT),
]
yy = 2.0
for t, d, col in flow:
    rect(s, rx + 0.3, yy, 0.16, 0.16, fill=col, shape=MSO_SHAPE.OVAL)
    text(s, rx + 0.6, yy - 0.05, rw - 0.9, 0.3,
         [[(t, 11.5, TEXT, True), ("   " + d, 10.5, MUTED, False)]])
    if col != CRIT:
        line(s, rx + 0.38, yy + 0.2, rx + 0.38, yy + 0.42, color=BORDER, w=1.0)
    yy += 0.47
footer(s, "06 · CM write-back")

# ============================================================ SLIDE 14: RADIO ENGINE
s = slide_base()
header(s, "Section 07", "Radio analytics engine: one scoring model")
text(s, 0.55, 1.28, 12.2, 0.3,
     [[("Every radio insight module emits the same normalized ", 12, MUTED, False),
       ("issue", 12, TEAL, True, MONO), (" — from core/radio/scoring.py:", 12, MUTED, False)]])
# code card
rect(s, 0.55, 1.68, 6.4, 3.85, fill=RGBColor(0x0A,0x10,0x18), line=BORDER, radius=0.03)
code = [
    ("def issue(module, category, title,", TEAL),
    ("          summary, score, cells, …):", TEAL),
    ("    severity = severity_from_score(score)", TEXT),
    ("    return {", MUTED),
    ("      'id': stable_id(module, category,", ACCENT),
    ("                      title, site_id, cells),", ACCENT),
    ("      'severity': severity,   # Crit→Info", TEXT),
    ("      'score': round(score, 2),", TEXT),
    ("      'cells': cells, 'site_id': …,", MUTED),
    ("      'evidence': {kpi, pre, post, delta},", GOLD),
    ("      'recommendation': …,", MUTED),
    ("      'source_url': '/capacity-hotspots' }", MUTED),
]
yy = 1.85
for t, c in code:
    text(s, 0.78, yy, 6.0, 0.3, [[(t, 11, c, False, MONO)]])
    yy += 0.305
# right explanation
text(s, 7.2, 1.68, 5.5, 0.3, [[("Why this matters", 14, TEAL, True)]])
bullets(s, 7.2, 2.06, 5.55, 3.0, [
    ("Uniform issue shape. ", "Filtering, sorting, severity roll-ups and the UI are written once and reused by every module."),
    ("Deterministic IDs. ", "stable_id = sha256(parts)[:16] — the same problem keeps the same id across runs, so dedup & tracking work."),
    ("Bounded, additive scores. ", "bounded_score sums capped signals → 0–100; thresholds map to Critical/High/Medium/Low/Info."),
    ("Evidence travels with the verdict. ", "pre/post/delta KPIs are attached, so “why” is auditable, not a black box."),
], size=12, gap=7)
# severity legend strip
sevs = [("Critical ≥85", CRIT), ("High ≥70", ORANGE), ("Medium ≥45", GOLD), ("Low >0", TEAL), ("Info 0", FAINT)]
cx = 0.55
for lab, col in sevs:
    w = 0.2 + 0.1 * len(lab)
    chip(s, cx, 5.68, w, lab, col, h=0.32, size=10)
    cx += w + 0.18
footer(s, "07 · Radio engine")

# ============================================================ SLIDE 15: INSIGHT CATALOG
s = slide_base()
header(s, "Section 07", "The insight catalog")
text(s, 0.55, 1.25, 12.2, 0.3,
     [[("Same engine, different recipes — each module scores a specific failure mode:", 12, MUTED, False)]])
cat = [
    ("Capacity Hotspots", "busy-hour PRB/utilization delta vs baseline → carrier/split advice", ACCENT),
    ("Sleeping Cells", "traffic/accessibility collapse vs history → outage-suspect cells", CRIT),
    ("Overshooting", "TA / distance vs footprint → down-tilt / power candidates", ORANGE),
    ("Layer Coverage Gaps", "missing LTE band per sector from coverage payload", PURPLE),
    ("Neighbor Quality", "relation health, missing/one-way neighbors, HO KPIs", TEAL),
    ("Sector Health", "per-sector KPI composite, all-cells variant", ACCENT),
    ("SON Analytics", "cross-cutting optimization insights & priorities", GOLD),
    ("Change Impact", "before/after KPI deltas around a config change", TEAL),
    ("RF Optimization", "workbench aggregating candidate actions", PURPLE),
    ("Radio Morning Report", "overnight roll-up of the above into one brief", ORANGE),
]
cw, rh = 6.05, 0.82
for i, (t, d, col) in enumerate(cat):
    cx = 0.55 + (i % 2) * (cw + 0.12)
    cy = 1.65 + (i // 2) * (rh + 0.06)
    rect(s, cx, cy, cw, rh, fill=PANEL, line=BORDER, radius=0.06)
    rect(s, cx, cy, 0.07, rh, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx + 0.24, cy + 0.09, cw - 0.4, 0.3, [[(t, 12, TEXT, True)]])
    text(s, cx + 0.24, cy + 0.42, cw - 0.4, 0.35, [[(d, 10, MUTED, False)]])
# compact inline demo pointer (full catalog fills the slide)
rect(s, 0.55, 6.28, 12.2, 0.4, fill=RGBColor(0x22,0x1C,0x0E), line=GOLD, line_w=1.0, radius=0.12)
text(s, 0.55, 6.28, 12.2, 0.4,
     [[("⟳  SWITCH TO TOOL:  ", 11, GOLD, True),
       ("open any two of these side-by-side and show the identical issue table + severity roll-up.", 10.5, MUTED, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, "07 · Radio engine")

# ============================================================ SLIDE 16: FRONTEND
s = slide_base()
header(s, "Section 08", "Frontend architecture")
cols = [
    ("Theme system", TEAL, [
        "body.dark-mode is the single switch; all CSS keys off it.",
        "Persisted in localStorage (primenet-theme); applied on every load.",
        "primenet:theme-change event syncs Chart.js instances live.",
        "Dark tokens (--dm-bg/panel/text…) in common.css; module CSS inherits.",
    ]),
    ("Shared shells", ACCENT, [
        "common.css / common.js — header, toggle, tables, buttons.",
        "radio_module.html — one filter shell for all radio modules.",
        "dashboard = constellation.css + constellation.js (always-dark deck).",
        "Cache-bust ?v=X.X bumped only on files actually changed.",
    ]),
    ("Rendering & maps", PURPLE, [
        "Server-rendered Jinja + progressive JS per module.",
        "Chart.js for KPI trends; theme-aware defaults.",
        "Leaflet maps — CSP allow-lists OSM & ArcGIS tile hosts.",
        "unpkg is the only external script/style origin permitted.",
    ]),
]
for i, (t, col, items) in enumerate(cols):
    cx = 0.55 + i * (4.05)
    rect(s, cx, 1.4, 3.9, 4.2, fill=PANEL, line=BORDER, radius=0.04)
    rect(s, cx, 1.4, 3.9, 0.5, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx + 0.22, 1.4, 3.6, 0.5, [[(t, 13, BG, True)]], anchor=MSO_ANCHOR.MIDDLE)
    yy = 2.1
    for it in items:
        text(s, cx + 0.22, yy, 3.5, 0.8, [[("▸ ", 11, col, True), (it, 11, MUTED, False)]], line_spacing=1.04)
        yy += 0.85
footer(s, "08 · Frontend")

# ============================================================ SLIDE 17: DEPLOYMENT
s = slide_base()
header(s, "Section 08", "Deployment & operations")
bullets(s, 0.55, 1.4, 6.05, 3.8, [
    ("Single container. ", "Dockerfile + gunicorn (GUNICORN_WORKERS/THREADS/TIMEOUT); docker-compose maps PRIMENET_HTTP_PORT."),
    ("Stateful volume. ", "Mount NCM_DATA_ROOT=/data — databases, raw, sync_downloads persist outside the image."),
    ("Config via env. ", ".env holds SFTP + CM/FM creds and tuning; nothing secret is committed."),
    ("Health probe. ", "/health & /api/health report locked / degraded / ok for LBs and orchestrators."),
    ("Hardening switches. ", "NCM_ENABLE_HSTS when HTTPS-only; NCM_CONTAINER disables dev-only Windows helpers."),
    ("Bootstrap. ", "First run seeds an admin from NCM_BOOTSTRAP_ADMIN_* only when the users table is empty."),
], size=12, gap=8)
# env card
rx, rw = 6.9, 5.85
rect(s, rx, 1.4, rw, 3.85, fill=RGBColor(0x0C,0x13,0x1D), line=BORDER, radius=0.04)
text(s, rx + 0.28, 1.55, rw - 0.5, 0.3, [[("KEY ENVIRONMENT", 11, MUTED, True)]])
env = [
    ("NCM_DATA_ROOT=/data", ACCENT),
    ("NCM_CONTAINER=1", FAINT),
    ("GUNICORN_WORKERS=2  THREADS=4", TEAL),
    ("FLASK_SECRET_KEY=…", GOLD),
    ("SESSION_LIFETIME_HOURS=2", FAINT),
    ("NOKIA_PM_HOST / HUAWEI_PM_HOST / METADATA_HOST", TEAL),
    ("NOKIA_CM_HOST / USER / PASSWORD", PURPLE),
    ("NOKIA_CM_MO_BATCH_SIZE=150", FAINT),
    ("NOKIA_CM_BATCH_DELAY_SEC=0.4", FAINT),
    ("NOKIA_FM_* (Keycloak/OAuth, Fault Mgmt)", CRIT),
    ("NCM_ENABLE_HSTS=1", GOLD),
]
yy = 1.95
for t, c in env:
    text(s, rx + 0.3, yy, rw - 0.5, 0.3, [[(t, 11, c, False, MONO)]])
    yy += 0.29
footer(s, "08 · Deployment")

# ============================================================ SLIDE 18: EXTENDING
s = slide_base()
header(s, "Section 08", "Extending PrimeNet: add a module")
steps = [
    ("1", "Scaffold the folder", "modules/<name>/ with routes.py, templates/, static/, logic.py — copy an existing module as the template.", ACCENT),
    ("2", "Define the blueprint", "<name>_bp = Blueprint(..., template_folder=\"templates\"); guard routes with login_required.", TEAL),
    ("3", "Register once", "Import + app.register_blueprint(<name>_bp) in app.py — the only global edit.", PURPLE),
    ("4", "Wire visibility", "Add the route to core/module_access.py with all / admin / admin_or_noc.", GOLD),
    ("5", "Reuse the shell", "Load common.css/js (+ radio_module.html for radio filters) so theme & chrome come free.", ORANGE),
    ("6", "Emit issues (if analytic)", "Return core/radio.scoring.issue() dicts to inherit filtering, severity & UI.", CRIT),
]
for i, (n, t, d, col) in enumerate(steps):
    cx = 0.55 + (i % 2) * 6.15
    cy = 1.45 + (i // 2) * 1.4
    rect(s, cx, cy, 5.95, 1.25, fill=PANEL, line=BORDER, radius=0.06)
    rect(s, cx, cy, 0.6, 1.25, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx, cy, 0.6, 1.25, [[(n, 26, BG, True, FONT_L)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 0.78, cy + 0.14, 5.0, 0.35, [[(t, 13, TEXT, True)]])
    text(s, cx + 0.78, cy + 0.5, 5.05, 0.7, [[(d, 10.5, MUTED, False)]], line_spacing=1.04)
footer(s, "08 · Extending")

# ============================================================ SLIDE 19: DEMO RUNBOOK
s = slide_base(divider=True)
constellation(s, n=40, seed=3, link=1.7)
rect(s, 0.55, 0.5, 0.09, 0.62, fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
text(s, 0.78, 0.46, 11.8, 0.3, [[("APPENDIX", 11, GOLD, True)]])
text(s, 0.77, 0.7, 12.0, 0.6, [[("End-to-end demo runbook", 25, TEXT, True, FONT_L)]])
text(s, 0.55, 1.35, 12.2, 0.3,
     [[("A single continuous click-path — narrate each subsystem as you pass through it:", 12, MUTED, False)]])
run = [
    ("Login & Dashboard", "Show the constellation deck & RBAC-filtered nav.", "Overview of modules by role"),
    ("Sync", "Trigger/observe a pull→load; then Network Health for freshness.", "ETL taxonomy in action"),
    ("Performance Explorer", "Query KPIs across vendor/RAT.", "Automatic column detection"),
    ("CM Extractor", "Extract a scope → open the Excel output.", "One UI, Nokia REST + Huawei MML"),
    ("Capacity Hotspots", "Sort by score; expand evidence on a row.", "The shared issue/scoring model"),
    ("RET Management", "Read tilt; walk the safe write flow (don't apply).", "Governed CM write-back"),
    ("Config History", "Show a recorded change.", "Auditability & rollback"),
]
yy = 1.85
for i, (t, d, why) in enumerate(run):
    rect(s, 0.55, yy, 12.2, 0.66, fill=RGBColor(0x14,0x1E,0x2B), line=BORDER, radius=0.05)
    rect(s, 0.55, yy, 0.5, 0.66, fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
    text(s, 0.55, yy, 0.5, 0.66, [[(str(i+1), 18, BG, True, FONT_L)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.2, yy, 3.1, 0.66, [[(t, 12.5, TEXT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.35, yy, 5.0, 0.66, [[(d, 11, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.45, yy, 3.15, 0.66, [[("↳ " + why, 10.5, TEAL, False, FONT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.735

# ============================================================ SLIDE 20: CLOSE
s = slide_base(divider=True)
constellation(s, n=64, seed=11, link=2.0)
text(s, 0.9, 2.2, 11.5, 0.4, [[("TAKEAWAYS", 14, ACCENT, True)]])
text(s, 0.85, 2.6, 11.6, 1.0, [[("One contract, 37 modules", 44, WHITE, True, FONT_L)]])
bullets(s, 0.9, 3.85, 11.5, 2.0, [
    ("A modular monolith ", "trades microservice overhead for a single, legible, single-container deploy."),
    ("Data flow is deterministic ", "— a path taxonomy + idempotent ETL make SQLite a feature, not a limitation."),
    ("Two vendors, one surface ", "— REST vs MML complexity is absorbed in core/, never in modules."),
    ("Analytics share one scoring model ", "— uniform issues make new insights cheap and auditable."),
], size=13, gap=8)
rect(s, 0.9, 6.15, 11.5, 0.7, fill=RGBColor(0x14,0x1E,0x2B), line=GOLD, line_w=1.0, radius=0.06)
text(s, 0.9, 6.15, 11.5, 0.7,
     [[("Questions?  ", 14, GOLD, True), ("Let's go back to the tool and dig into whatever you want to see.", 13, TEXT, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

out = "/home/user/PrimeNet/PrimeNet_Engineering_DeepDive.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
