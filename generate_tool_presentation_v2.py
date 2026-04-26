from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent
LOGO_PATH = PROJECT_ROOT / "static" / "images" / "favicon.png"

# PrimeNet-like palette
BRAND_PRIMARY = RGBColor(127, 166, 194)
BRAND_DARK = RGBColor(44, 62, 80)
BRAND_ACCENT = RGBColor(52, 152, 219)
TEXT_LIGHT = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(44, 62, 80)


def _paint_title(slide, title: str, subtitle: str | None = None):
    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_PRIMARY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.8), Inches(0.8))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = TEXT_LIGHT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.0), Inches(0.5))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = BRAND_DARK

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(11.9), Inches(0.12), height=Inches(0.95))


def _bullets_box(slide, left, top, width, height, bullets):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_DARK


def add_brand_title(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _paint_title(slide, title)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.8), Inches(2.0))
    tf = sub.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BRAND_DARK
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = "Prepared by: Malek Mohammad"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(90, 105, 120)

    return slide


def add_bullets_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_title(slide, title, subtitle)
    _bullets_box(slide, Inches(0.8), Inches(2.0), Inches(12.0), Inches(4.7), bullets)
    return slide


def add_screenshots_placeholder(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_title(slide, "Live Screenshots", "Replace placeholders with current UI captures")

    slots = [
        ("Dashboard", 0.7, 2.0),
        ("Performance Analytics", 4.6, 2.0),
        ("Network Map / Neighbor Analysis", 8.5, 2.0),
    ]
    for label, x, y in slots:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.9), Inches(3.6))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(238, 244, 250)
        rect.line.color.rgb = BRAND_PRIMARY
        t = rect.text_frame
        t.clear()
        p = t.paragraphs[0]
        p.text = f"{label}\n\nInsert screenshot here"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = BRAND_DARK
    return slide


def build_braned_full(path_out: Path):
    prs = Presentation()
    add_brand_title(prs, "PrimeNet Platform", "Network Performance & Configuration Tool")
    add_bullets_slide(
        prs,
        "What Problem We Solve",
        [
            "Fragmented network tools and manual spreadsheet workflows",
            "Slow root-cause analysis across KPI, topology, and config data",
            "Inconsistent reporting across operations and optimization teams",
        ],
    )
    add_bullets_slide(
        prs,
        "Core Modules",
        [
            "Performance Analytics: KPI trends, DOD, PM table exploration",
            "Network Map: site view, sector wedges, search and filters",
            "Neighbor Analysis: HO links, attempts/failures, source-target insights",
            "Reports and exports for operational and management use",
        ],
    )
    add_bullets_slide(
        prs,
        "Workflow",
        [
            "Automated SFTP pull + normalization into data stores",
            "User selects scope (technology/area/cluster/object)",
            "Query KPIs and visualize trends with contextual map overlays",
            "Export findings for decisions and executive reporting",
        ],
    )
    add_screenshots_placeholder(prs)
    add_bullets_slide(
        prs,
        "Business Value",
        [
            "Faster troubleshooting and reduced time to isolate issues",
            "More reliable daily monitoring and trend visibility",
            "Less manual effort and better data consistency",
            "Single platform for NOC, optimization, and management teams",
        ],
    )
    add_bullets_slide(
        prs,
        "Roadmap",
        [
            "Data quality guards for neighbor and KPI feed integrity",
            "Scalable map rendering for dense nationwide views",
            "Role-based dashboards and scheduled executive summaries",
        ],
    )
    prs.save(str(path_out))


def build_executive_cut(path_out: Path):
    prs = Presentation()
    add_brand_title(prs, "PrimeNet Executive Brief", "5-slide management summary")
    add_bullets_slide(
        prs,
        "1) Platform in One Minute",
        [
            "Unified tool for network performance, topology, and configuration",
            "Built to shorten troubleshooting cycles and improve decision speed",
        ],
    )
    add_bullets_slide(
        prs,
        "2) Key Capabilities",
        [
            "KPI analytics (trend + day-over-day)",
            "Map-based network context with sectors and neighbors",
            "Operational reports and export-ready outputs",
        ],
    )
    add_bullets_slide(
        prs,
        "3) Operational Impact",
        [
            "Faster issue isolation and clearer ownership",
            "Standardized monitoring workflow across teams",
            "Reduced dependence on manual data preparation",
        ],
    )
    add_bullets_slide(
        prs,
        "4) Next Steps",
        [
            "Adopt role-specific dashboards",
            "Expand automation and data quality validation",
            "Institutionalize monthly KPI business reviews",
        ],
    )
    prs.save(str(path_out))


if __name__ == "__main__":
    build_braned_full(PROJECT_ROOT / "PrimeNet_Tool_Overview_Branded.pptx")
    build_executive_cut(PROJECT_ROOT / "PrimeNet_Tool_Overview_Executive_5slides.pptx")
