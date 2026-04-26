from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
    return slide


def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = title

    left = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(6.2), Inches(5.2))
    right = slide.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(6.2), Inches(5.2))

    for box, box_title, points in ((left, left_title, left_points), (right, right_title, right_points)):
        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = box_title
        p0.font.bold = True
        p0.font.size = Pt(26)
        p0.font.color.rgb = RGBColor(36, 67, 98)
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.level = 0
            p.font.size = Pt(20)

    return slide


def add_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = "Roadmap & Next Steps"

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12.2), Inches(5.4))
    tf = box.text_frame
    tf.clear()

    title = tf.paragraphs[0]
    title.text = "Planned Enhancements"
    title.font.bold = True
    title.font.size = Pt(26)
    title.font.color.rgb = RGBColor(36, 67, 98)

    items = [
        "Automated quality checks for neighbor exports before loading",
        "Faster map rendering with marker clustering for dense views",
        "More executive reporting templates and scheduled exports",
        "Role-based dashboards and team-specific KPI views",
        "Operational alerts for KPI threshold breaches",
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)

    return slide


def build_presentation(output_path):
    prs = Presentation()

    add_title_slide(
        prs,
        "PrimeNet Platform Overview",
        "Network Performance & Configuration Tool\nPrepared by: Malek Mohammad",
    )

    add_bullets_slide(
        prs,
        "Why This Tool Exists",
        [
            "Unifies network performance, topology, and configuration workflows",
            "Reduces manual Excel handling and fragmented analysis steps",
            "Accelerates troubleshooting with map + KPI context in one place",
            "Supports daily operations, optimization, and management reporting",
        ],
    )

    add_bullets_slide(
        prs,
        "Core Modules",
        [
            "Performance Analytics: KPI trends, DOD view, PM database exploration",
            "Network Map: site view, sectors, technology filtering, search and zoom",
            "Neighbor Analysis: handover links, source/target diagnostics, failure views",
            "Reports: reusable report definitions and export-ready outputs",
            "Configuration Utilities: parameter dictionary, XML tools, task scheduler",
        ],
    )

    add_two_column_slide(
        prs,
        "Operational Workflow",
        "Data Ingestion",
        [
            "SFTP pulls for Nokia, Huawei, metadata and neighbor datasets",
            "Automated loaders normalize files into SQLite/PostgreSQL stores",
            "Retention policies keep data relevant and storage controlled",
        ],
        "User Analysis Flow",
        [
            "Select scope (technology, area, cluster, object)",
            "Query PM data and visualize KPI charts",
            "Use map and neighbor lines to isolate issues quickly",
            "Export evidence for reports and stakeholder updates",
        ],
    )

    add_bullets_slide(
        prs,
        "Business Value",
        [
            "Faster root-cause analysis across radio technologies",
            "Improved visibility into trends, anomalies, and network quality",
            "Higher consistency in reporting and operational decisions",
            "Lower manual effort for NOC and optimization teams",
        ],
    )

    add_roadmap_slide(prs)

    add_bullets_slide(
        prs,
        "Thank You",
        [
            "Questions & discussion",
            "Next: tailor this deck with your KPIs and live screenshots",
        ],
    )

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("PrimeNet_Tool_Overview.pptx")
