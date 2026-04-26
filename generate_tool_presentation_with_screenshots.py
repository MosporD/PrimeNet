from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent
LOGO_PATH = PROJECT_ROOT / "static" / "images" / "favicon.png"

ASSETS = [
    # title, image path (from provided screenshots)
    (
        "Login Page",
        Path(r"C:\Users\malek.mohammad\.cursor\projects\c-Users-malek-mohammad-Project-Cursor-version-Project\assets\c__Users_malek.mohammad_AppData_Roaming_Cursor_User_workspaceStorage_ca470328415ae8bd6f846c68e50bb8e9_images_image-1a849dc5-c1b2-4b0a-a270-534df1965967.png"),
    ),
    (
        "Performance Analytics",
        Path(r"C:\Users\malek.mohammad\.cursor\projects\c-Users-malek-mohammad-Project-Cursor-version-Project\assets\c__Users_malek.mohammad_AppData_Roaming_Cursor_User_workspaceStorage_ca470328415ae8bd6f846c68e50bb8e9_images_image-12327826-4f8f-4b4d-b3de-0af607e20fae.png"),
    ),
    (
        "Network Map",
        Path(r"C:\Users\malek.mohammad\.cursor\projects\c-Users-malek-mohammad-Project-Cursor-version-Project\assets\c__Users_malek.mohammad_AppData_Roaming_Cursor_User_workspaceStorage_ca470328415ae8bd6f846c68e50bb8e9_images_image-495a68ef-3367-44dc-8798-4fb583e3f3b8.png"),
    ),
    (
        "Neighbor Analysis",
        Path(r"C:\Users\malek.mohammad\.cursor\projects\c-Users-malek-mohammad-Project-Cursor-version-Project\assets\c__Users_malek.mohammad_AppData_Roaming_Cursor_User_workspaceStorage_ca470328415ae8bd6f846c68e50bb8e9_images_image-72c0b6b1-6924-4576-b9db-dc6ba3baacbb.png"),
    ),
    (
        "Femto PM",
        Path(r"C:\Users\malek.mohammad\.cursor\projects\c-Users-malek-mohammad-Project-Cursor-version-Project\assets\c__Users_malek.mohammad_AppData_Roaming_Cursor_User_workspaceStorage_ca470328415ae8bd6f846c68e50bb8e9_images_image-b8a26ee5-ff26-4452-952e-3fb2c7ac0dc6.png"),
    ),
    (
        "Configuration Task Scheduler",
        Path(r"C:\Users\malek.mohammad\.cursor\projects\c-Users-malek-mohammad-Project-Cursor-version-Project\assets\c__Users_malek.mohammad_AppData_Roaming_Cursor_User_workspaceStorage_ca470328415ae8bd6f846c68e50bb8e9_images_image-5b016b8a-1197-4107-a35d-31a9a8a404c3.png"),
    ),
    (
        "Parameter Dictionary",
        Path(r"C:\Users\malek.mohammad\.cursor\projects\c-Users-malek-mohammad-Project-Cursor-version-Project\assets\c__Users_malek.mohammad_AppData_Roaming_Cursor_User_workspaceStorage_ca470328415ae8bd6f846c68e50bb8e9_images_image-17395fa4-bca7-41ff-a268-177b3652e128.png"),
    ),
    (
        "Dashboard",
        Path(r"C:\Users\malek.mohammad\.cursor\projects\c-Users-malek-mohammad-Project-Cursor-version-Project\assets\c__Users_malek.mohammad_AppData_Roaming_Cursor_User_workspaceStorage_ca470328415ae8bd6f846c68e50bb8e9_images_image-c15c2850-4a4e-4af2-b019-60189f85ac21.png"),
    ),
]

BRAND_PRIMARY = RGBColor(127, 166, 194)
BRAND_DARK = RGBColor(44, 62, 80)


def _paint_header(slide, title: str):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_PRIMARY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(10.8), Inches(0.6))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(12.0), Inches(0.08), height=Inches(0.84))


def _add_image_fit(slide, image_path: Path, left, top, width, height):
    with Image.open(image_path) as img:
        iw, ih = img.size
    frame_w = width.inches
    frame_h = height.inches
    img_ratio = iw / ih
    frame_ratio = frame_w / frame_h

    if img_ratio > frame_ratio:
        # image wider than frame -> fit width
        draw_w = frame_w
        draw_h = frame_w / img_ratio
    else:
        # image taller than frame -> fit height
        draw_h = frame_h
        draw_w = frame_h * img_ratio

    dx = (frame_w - draw_w) / 2
    dy = (frame_h - draw_h) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(left.inches + dx),
        Inches(top.inches + dy),
        width=Inches(draw_w),
        height=Inches(draw_h),
    )


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_header(slide, "PrimeNet")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(2.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Platform Walkthrough"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = BRAND_DARK

    p2 = tf.add_paragraph()
    p2.text = "Network Performance & Configuration Tool"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(90, 105, 120)


def add_screen_slide(prs, title: str, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_header(slide, title)

    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.15), Inches(12.43), Inches(5.95))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(244, 247, 251)
    frame.line.color.rgb = RGBColor(196, 214, 232)
    frame.line.width = Pt(1.2)

    if image_path.exists():
        _add_image_fit(slide, image_path, Inches(0.55), Inches(1.25), Inches(12.23), Inches(5.75))
    else:
        tx = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(10.5), Inches(0.8))
        tf = tx.text_frame
        tf.text = f"Screenshot missing: {image_path.name}"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RGBColor(160, 60, 60)


def add_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_header(slide, "Thank You")
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BRAND_DARK

    p2 = tf.add_paragraph()
    p2.text = "PrimeNet live demonstration and roadmap available."
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(90, 105, 120)


def build():
    prs = Presentation()
    add_cover(prs)
    for title, image in ASSETS:
        add_screen_slide(prs, title, image)
    add_closing(prs)
    out = PROJECT_ROOT / "PrimeNet_Tool_Overview_With_Screenshots.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
