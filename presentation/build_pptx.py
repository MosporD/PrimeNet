"""Generate PrimeNet_Presentation.pptx — a native, editable PowerPoint deck.

Mirrors the reveal.js HTML deck: dark navy theme, generated illustrations,
capability pillars, architecture + pipeline diagrams, feature catalog, roadmap.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")

# ---------- palette ----------
NAVY      = RGBColor(0x0A, 0x13, 0x26)
NAVY_DK   = RGBColor(0x06, 0x0C, 0x19)
CARD      = RGBColor(0x14, 0x21, 0x38)
CARD_HI   = RGBColor(0x18, 0x2A, 0x48)
CYAN      = RGBColor(0x22, 0xD3, 0xEE)
TEAL      = RGBColor(0x2D, 0xD4, 0xBF)
BLUE      = RGBColor(0x3B, 0x82, 0xF6)
INK       = RGBColor(0xEA, 0xF2, 0xFF)
MUTED     = RGBColor(0x9F, 0xB3, 0xD1)
LINE      = RGBColor(0x2A, 0x3D, 0x63)
NOKIA     = RGBColor(0x38, 0xBD, 0xF8)
HUAWEI    = RGBColor(0xFB, 0x71, 0x85)
ROADBLUE  = RGBColor(0x60, 0xA5, 0xFA)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------- low-level helpers ----------
def _set_alpha(fill_elem_color, alpha_pct):
    """Add an <a:alpha> child to a color element (0-100)."""
    srgb = fill_elem_color._xFill.find(qn('a:solidFill'))
    # not used directly; handled in solid_alpha


def solid(shape, color, alpha=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if alpha is not None:
        sf = shape.fill.fore_color._xFill  # a:solidFill
        srgb = sf.find(qn('a:srgbClr'))
        a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha * 1000))})
        srgb.append(a)
    shape.line.fill.background()


def no_line(shape):
    shape.line.fill.background()


def border(shape, color=LINE, w=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def add_slide(bg=NAVY):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    solid(r, bg)
    r.shadow.inherit = False
    # top accent hairline
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Pt(4))
    solid(bar, CYAN)
    bar.shadow.inherit = False
    return s


def box(slide, l, t, w, h, shape=MSO_SHAPE.ROUNDED_RECTANGLE, fill=CARD,
        line_color=LINE, line_w=0.75):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    solid(sp, fill)
    if line_color is None:
        no_line(sp)
    else:
        border(sp, line_color, line_w)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    return sp


def _set_run(r, text, size, color, bold=False, italic=False, font=FONT):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, space_after=2, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (text,size,color,bold) tuples."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for seg in para:
            txt, size, color, bold = (seg + (False,))[:4] if len(seg) == 3 else seg
            run = p.add_run()
            _set_run(run, txt, size, color, bold)
    return tb


def eyebrow(slide, label, l=Inches(0.6), t=Inches(0.45)):
    w = Inches(0.18 + 0.085 * len(label))
    chip = box(slide, l, t, w, Inches(0.34), shape=MSO_SHAPE.ROUNDED_RECTANGLE,
               fill=RGBColor(0x10, 0x2A, 0x40), line_color=LINE)
    text(slide, l, t, w, Inches(0.34),
         [[(label.upper(), 10.5, CYAN, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return chip


def title(slide, txt, l=Inches(0.6), t=Inches(0.92), w=Inches(12.1), size=30):
    text(slide, l, t, w, Inches(0.9), [[(txt, size, INK, True)]], line_spacing=1.05)


def card(slide, l, t, w, h, icon, head, body, accent=CYAN, head_size=14,
         body_size=10.5, fill=CARD, line_color=LINE):
    c = box(slide, l, t, w, h, fill=fill, line_color=line_color)
    pad = Inches(0.18)
    runs = []
    if icon:
        runs.append([(icon + "  ", 15, accent, True), (head, head_size, INK, True)])
    else:
        runs.append([(head, head_size, INK, True)])
    runs.append([(body, body_size, MUTED, False)])
    tb = text(slide, l + pad, t + Inches(0.14), w - 2 * pad, h - Inches(0.24), runs,
              line_spacing=1.06, space_after=4)
    return c


def chip(slide, l, t, label, color=INK, bcolor=LINE, fill=RGBColor(0x10, 0x1B, 0x30)):
    w = Inches(0.22 + 0.10 * len(label))
    c = box(slide, l, t, w, Inches(0.36), fill=fill, line_color=bcolor)
    text(slide, l, t, w, Inches(0.36), [[(label, 10.5, color, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return l + w + Inches(0.12)


def img_box(slide, path, l, t, w, h):
    """Place an image cropped to fill the box (cover)."""
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 1024, 683
    box_ratio = w / h
    img_ratio = iw / ih
    pic = slide.shapes.add_picture(path, l, t, w, h)
    if img_ratio > box_ratio:
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - img_ratio / box_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    pic.line.color.rgb = LINE
    pic.line.width = Pt(1)
    return pic


def arrow(slide, l, t, w=Inches(0.5), h=Inches(0.4), color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    solid(a, color)
    a.shadow.inherit = False
    try:
        a.adjustments[0] = 0.5
        a.adjustments[1] = 0.55
    except Exception:
        pass
    return a


def footer(slide, txt):
    text(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35),
         [[(txt, 9, RGBColor(0x6B, 0x7D, 0x9C), False)]], align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
solid(bg, NAVY_DK); bg.shadow.inherit = False
hero = os.path.join(ASSETS, "hero.png")
if os.path.exists(hero):
    img_box(s, hero, 0, 0, EMU_W, EMU_H)
veil = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
solid(veil, NAVY_DK, alpha=72); veil.shadow.inherit = False
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Pt(4))
solid(bar, CYAN); bar.shadow.inherit = False

eyebrow(s, "Internal Operations Platform · Multi-Vendor RAN",
        l=Inches(3.4), t=Inches(2.05))
text(s, Inches(0.6), Inches(2.55), Inches(12.1), Inches(1.4),
     [[("PrimeNet", 72, CYAN, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(3.95), Inches(12.1), Inches(0.6),
     [[("Radio Network Performance & Configuration Platform", 22, INK, True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.6), Inches(4.6), Inches(10.1), Inches(0.9),
     [[("One governed workspace for performance intelligence, configuration", 13, MUTED, False)],
      [("discipline, and daily tool access across Nokia & Huawei networks.", 13, MUTED, False)]],
     align=PP_ALIGN.CENTER, line_spacing=1.3)
# chips row centered
labels = [("Nokia", NOKIA), ("Huawei", HUAWEI), ("2G · 3G · 4G · 5G", INK),
          ("Flask", INK), ("SQLite", INK), ("APScheduler", INK)]
total = sum(Inches(0.22 + 0.10 * len(l)) + Inches(0.12) for l, _ in labels) - Inches(0.12)
x = int((EMU_W - total) / 2)
for lab, col in labels:
    x = chip(s, x, Inches(5.55), lab, color=col,
             bcolor=(col if col in (NOKIA, HUAWEI) else LINE))
footer(s, "Zain Jordan RAN context  ·  Executive & stakeholder briefing  ·  June 2026")


# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
s = add_slide()
eyebrow(s, "The Problem")
title(s, "Radio operations live in too many disconnected systems")
bullets = [
    ("10+ separate tools", " — vendor OSS, assurance portals, drive-test, site DBs, spreadsheets."),
    ("Nokia & Huawei data live apart", " — inconsistent KPI definitions slow decisions."),
    ("Manual exports & merges", " — analysts spend hours reconciling before analysis."),
    ("No config discipline", " — nowhere to compare files, track tasks, retain history."),
    ("No unified picture", " — leadership lacks one view of sites by tech & vendor."),
]
runs = []
for b, rest in bullets:
    runs.append([("▸  ", 13, CYAN, True), (b, 12.5, INK, True), (rest, 12.5, MUTED, False)])
text(s, Inches(0.6), Inches(1.95), Inches(6.7), Inches(4.5), runs, line_spacing=1.25, space_after=12)

cards2 = [("🧭", "Context switching", "Every incident means logging into multiple systems."),
          ("📑", "Spreadsheet sprawl", "Ad-hoc merges that don't agree with each other."),
          ("⏱️", "Slow to insight", "Hours between data drop and decision."),
          ("⚠️", "Config drift", "Undetected changes, no audit trail.")]
gx, gy, cw, ch, gap = Inches(7.55), Inches(1.95), Inches(2.65), Inches(2.1), Inches(0.25)
for i, (ic, h, b) in enumerate(cards2):
    cx = gx + (i % 2) * (cw + gap)
    cy = gy + (i // 2) * (ch + gap)
    card(s, cx, cy, cw, ch, ic, h, b)


# ============================================================
# SLIDE 3 — THE SOLUTION
# ============================================================
s = add_slide()
eyebrow(s, "The Solution")
title(s, "One secure web platform that unifies the workflow")
pipe = os.path.join(ASSETS, "pipeline.png")
if os.path.exists(pipe):
    img_box(s, pipe, Inches(0.6), Inches(2.0), Inches(5.7), Inches(4.3))
sol = [
    ("Automates collection", " of PM, grouping, neighbor, metadata & small-cell data on hourly/daily schedules."),
    ("Stores & serves", " it through one dashboard with consistent KPI definitions across vendors."),
    ("Supports config discipline", " — lookup, convert, compare, schedule, version history."),
    ("Acts as an operations hub", " — live site metrics plus one-click launch of existing tools."),
]
runs = []
for b, rest in sol:
    runs.append([("▸  ", 13, CYAN, True), (b, 12.5, INK, True), (rest, 12.5, MUTED, False)])
text(s, Inches(6.7), Inches(2.05), Inches(6.0), Inches(3.4), runs, line_spacing=1.22, space_after=12)
note = box(s, Inches(6.7), Inches(5.5), Inches(6.0), Inches(0.8), fill=RGBColor(0x0E, 0x22, 0x33), line_color=TEAL)
text(s, Inches(6.9), Inches(5.5), Inches(5.6), Inches(0.8),
     [[("PrimeNet complements vendor OSS — it does not replace it.", 13, RGBColor(0x9F, 0xE9, 0xFF), True)]],
     anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 4 — AT A GLANCE
# ============================================================
s = add_slide(NAVY_DK)
eyebrow(s, "At a Glance")
title(s, "A modular Flask monolith built for scale")
stats = [("19", "Feature blueprints"), ("2", "Vendors unified\n(Nokia · Huawei)"),
         ("4", "Technologies\n2G / 3G / 4G / 5G"), ("7×24", "Automated hourly\n+ daily ingest"),
         ("8+", "SQLite data stores\n(WAL, read-optimized)"), ("4", "Access roles\nuser → admin"),
         ("SFTP", "Vendor feeds\n+ REST / MML APIs"), ("KML", "+ Excel / XML\nexport formats")]
cw, ch, gap = Inches(2.85), Inches(1.85), Inches(0.25)
x0 = Inches(0.7)
for i, (num, lbl) in enumerate(stats):
    cx = x0 + (i % 4) * (cw + gap)
    cy = Inches(2.2) + (i // 4) * (ch + gap)
    box(s, cx, cy, cw, ch, fill=CARD, line_color=LINE)
    text(s, cx, cy + Inches(0.22), cw, Inches(0.8), [[(num, 40, CYAN, True)]], align=PP_ALIGN.CENTER)
    lines = lbl.split("\n")
    text(s, cx, cy + Inches(1.08), cw, Inches(0.7),
         [[(ln, 11, MUTED, False)] for ln in lines], align=PP_ALIGN.CENTER, line_spacing=1.05)


# ============================================================
# SLIDE 5 — THREE PILLARS
# ============================================================
s = add_slide()
eyebrow(s, "Capability Pillars")
title(s, "Three pillars, one workspace")
pil = os.path.join(ASSETS, "pillars.png")
if os.path.exists(pil):
    img_box(s, pil, Inches(0.6), Inches(2.0), Inches(4.7), Inches(4.3))
pillars = [
    (CYAN, "📈", "Performance", "KPI analytics & trends · network map & neighbors · reports · conflict / PCI views · Femto monitoring."),
    (TEAL, "⚙️", "Configuration", "Parameter dictionary · XML ⇄ Excel · NE compare · config tasks · version history."),
    (BLUE, "🛰️", "Operations Hub", "Live site dashboard by tech & vendor · external tool launcher · drive tests · admin & data health."),
]
px, pw, ph, pgap = Inches(5.6), Inches(7.1), Inches(1.25), Inches(0.2)
for i, (ac, ic, h, b) in enumerate(pillars):
    cy = Inches(2.0) + i * (ph + pgap)
    card(s, px, cy, pw, ph, ic, h, b, accent=ac, head_size=15, body_size=11)


# ============================================================
# SLIDE 6 — ARCHITECTURE
# ============================================================
s = add_slide(NAVY_DK)
eyebrow(s, "System Architecture")
title(s, "Everything at a glance")

def diagram_box(slide, l, t, w, h, head, sub, fill, lc):
    b = box(slide, l, t, w, h, fill=fill, line_color=lc, line_w=1.25)
    runs = [[(head, 12.5, INK, True)]]
    if sub:
        runs.append([(sub, 9.5, MUTED, False)])
    text(slide, l + Inches(0.1), t, w - Inches(0.2), h, runs,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    return b

# columns: External -> Pipeline -> PrimeNet core -> SQLite
colY = Inches(2.0)
diagram_box(s, Inches(0.55), colY, Inches(2.7), Inches(1.0), "External Sources",
            "SFTP feeds · CM/FM/PM REST & MML · uploads", RGBColor(0x24, 0x13, 0x27), HUAWEI)
diagram_box(s, Inches(0.55), Inches(3.4), Inches(2.7), Inches(1.0), "Ingestion Pipeline",
            "Pull orchestrators → loaders + processors", RGBColor(0x0E, 0x22, 0x33), TEAL)

diagram_box(s, Inches(3.75), Inches(1.85), Inches(3.1), Inches(0.95), "app.py",
            "before / after hooks", RGBColor(0x0C, 0x20, 0x36), CYAN)
diagram_box(s, Inches(3.75), Inches(2.95), Inches(3.1), Inches(0.95), "19 Feature Blueprints",
            "performance · map · config · admin", RGBColor(0x0C, 0x20, 0x36), CYAN)
diagram_box(s, Inches(3.75), Inches(4.05), Inches(3.1), Inches(0.95), "APScheduler + ncm_core",
            "jobs · XML / Excel converters", RGBColor(0x0C, 0x20, 0x36), CYAN)

diagram_box(s, Inches(7.35), Inches(1.85), Inches(2.55), Inches(0.78), "PM cells DBs", "", RGBColor(0x10, 0x18, 0x2E), ROADBLUE)
diagram_box(s, Inches(7.35), Inches(2.78), Inches(2.55), Inches(0.78), "metadata.db", "", RGBColor(0x10, 0x18, 0x2E), ROADBLUE)
diagram_box(s, Inches(7.35), Inches(3.71), Inches(2.55), Inches(0.78), "neighbors / groups", "", RGBColor(0x10, 0x18, 0x2E), ROADBLUE)
diagram_box(s, Inches(7.35), Inches(4.64), Inches(2.55), Inches(0.78), "ncm_users.db + sync_log", "", RGBColor(0x10, 0x18, 0x2E), ROADBLUE)

# browser + label
diagram_box(s, Inches(10.2), Inches(2.95), Inches(2.45), Inches(0.95), "🌐 Web Browser",
            "dashboards & tools", RGBColor(0x0B, 0x25, 0x40), NOKIA)

arrow(s, Inches(3.3), Inches(2.95), Inches(0.42), Inches(0.4))   # ext/pipe -> core
arrow(s, Inches(6.9), Inches(2.95), Inches(0.42), Inches(0.4))   # core -> db
arrow(s, Inches(9.95), Inches(3.3), Inches(0.42), Inches(0.4), color=NOKIA)  # db -> browser
text(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.7),
     [[("Flask · SQLite-only · APScheduler background jobs · cookie-session auth · read-optimized for dashboards while writes happen.", 11.5, MUTED, False)]],
     align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — DATA PIPELINE
# ============================================================
s = add_slide(NAVY_DK)
eyebrow(s, "Data Ingestion")
title(s, "From vendor file drop to analyst view — automatically")
steps = [
    ("⏱️ Triggers", "watcher · cron · admin", RGBColor(0x24, 0x13, 0x27), HUAWEI),
    ("Pull", "SFTP download", RGBColor(0x0E, 0x22, 0x33), TEAL),
    ("Staging", "sync_downloads · raw", RGBColor(0x10, 0x18, 0x2E), ROADBLUE),
    ("Load", "parse CSV/XLSX/ZIP", RGBColor(0x0C, 0x20, 0x36), CYAN),
    ("Processors", "upsert · retention", RGBColor(0x0C, 0x20, 0x36), CYAN),
    ("SQLite + sync_log", "→ Feature APIs → UI", RGBColor(0x10, 0x18, 0x2E), ROADBLUE),
]
bw, bh = Inches(1.78), Inches(1.4)
x = Inches(0.5)
y = Inches(2.7)
for i, (h, sub, fill, lc) in enumerate(steps):
    b = box(s, x, y, bw, bh, fill=fill, line_color=lc, line_w=1.25)
    text(s, x + Inches(0.08), y, bw - Inches(0.16), bh,
         [[(h, 12.5, INK, True)], [(sub, 9.5, MUTED, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    if i < len(steps) - 1:
        arrow(s, x + bw + Inches(0.02), y + Inches(0.5), Inches(0.32), Inches(0.4))
    x += bw + Inches(0.36)

modes = [
    ("Watcher-primary", "poll signatures, pull+load on change (default)"),
    ("Legacy-periodic", "hourly + daily full sync"),
    ("Manual-only", "admin API triggers"),
]
mx = Inches(0.7)
for h, b in modes:
    c = box(s, mx, Inches(4.85), Inches(3.85), Inches(1.0), fill=CARD, line_color=LINE)
    text(s, mx + Inches(0.18), Inches(4.85), Inches(3.5), Inches(1.0),
         [[(h, 12.5, RGBColor(0x9F, 0xE9, 0xFF), True)], [(b, 10.5, MUTED, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    mx += Inches(4.05)


# ============================================================
# SLIDE 8 — PERFORMANCE
# ============================================================
s = add_slide()
eyebrow(s, "Pillar 1 · Performance & Analytics")
title(s, "Hunt degradation. Compare vendors. Map the network.")
perf = [
    ("📊", "KPI Analytics", "Cell & group trends, dynamic columns, saved filter profiles — the largest engine in the app."),
    ("🗺️", "Network Map", "Geospatial sites, sector wedges, per-cell KPIs, search by cell code, KML / CSV export."),
    ("🔗", "Neighbor Analysis", "Neighbor relation lines & per-cell summaries from dedicated neighbor KPI databases."),
    ("📡", "Conflict / PCI Map", "Distance / bearing / sector-alignment risk scoring with strict→relaxed profiles."),
    ("🏠", "Femto PM", "Small-cell device catalog, KPI columns, and formula-driven trend evaluation."),
    ("📋", "Reports", "Generate, download & archive operational reports with metadata + elevation helpers."),
]
cw, ch, gx, gy, gap = Inches(3.95), Inches(1.55), Inches(0.6), Inches(1.95), Inches(0.22)
for i, (ic, h, b) in enumerate(perf):
    cx = gx + (i % 3) * (cw + gap)
    cy = gy + (i // 3) * (ch + gap)
    card(s, cx, cy, cw, ch, ic, h, b, body_size=10)
endpoints = "/api/performance/cell/trend   ·   /api/map/cells/wedge-data   ·   /api/network-map/neighbors/lines   ·   /api/conflict-map/export-kml"
ep = box(s, Inches(0.6), Inches(5.55), Inches(12.13), Inches(0.55), fill=RGBColor(0x0B, 0x1B, 0x2E), line_color=LINE)
text(s, Inches(0.6), Inches(5.55), Inches(12.13), Inches(0.55),
     [[(endpoints, 11, RGBColor(0x9F, 0xE9, 0xFF), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 9 — CONFIGURATION
# ============================================================
s = add_slide()
eyebrow(s, "Pillar 2 · Configuration Tooling")
title(s, "Safer changes with structured, auditable workflows")
tools = [
    ("XML Parser", "XML → Excel", CYAN),
    ("Excel Generator", "Excel → XML", CYAN),
    ("NE Comparison", "diff 2 NE XMLs", CYAN),
    ("Config History", "versions · diff · download", TEAL),
    ("Task Scheduler", "attachments · status · results", TEAL),
    ("Parameter Dictionary", "MO lookup + search", TEAL),
]
cw, ch, gx, gy, gap = Inches(3.95), Inches(1.25), Inches(0.6), Inches(2.0), Inches(0.22)
for i, (h, b, ac) in enumerate(tools):
    cx = gx + (i % 3) * (cw + gap)
    cy = gy + (i // 3) * (ch + gap)
    fill = RGBColor(0x0C, 0x20, 0x36) if i < 3 else RGBColor(0x0E, 0x22, 0x33)
    card(s, cx, cy, cw, ch, "", h, b, accent=ac, head_size=14, body_size=11, fill=fill, line_color=ac)
text(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(0.5),
     [[("Top row powered by ", 12, MUTED, False), ("ncm_core", 12, CYAN, True),
       (" converters · bottom row is stateful, stored in the app database.", 12, MUTED, False)]],
     align=PP_ALIGN.CENTER)
chips = [("CONVERT", "XML ⇄ Excel both directions"), ("COMPARE", "before / after NE drift"),
         ("TRACK", "version & schedule with traceability")]
cx = Inches(0.7)
for tag, desc in chips:
    c = box(s, cx, Inches(5.6), Inches(3.85), Inches(0.95), fill=CARD, line_color=LINE)
    text(s, cx + Inches(0.18), Inches(5.6), Inches(3.5), Inches(0.95),
         [[(tag, 10.5, TEAL, True)], [(desc, 11, INK, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    cx += Inches(4.05)


# ============================================================
# SLIDE 10 — OPERATIONS HUB
# ============================================================
s = add_slide()
eyebrow(s, "Pillar 3 · Operations Hub")
title(s, "The single front door for daily operations")
ops = [
    ("Live site dashboard", " — operational sites counted by RAT and vendor from active metadata."),
    ("External tool launcher", " — links to PRS, U2020, NetAct, NetChart, Nemo, TEMS."),
    ("Drive-test viewer", " — ingest GPX / NMFS field artifacts into map-ready traces."),
    ("Network management", " — browse sites & cells, run PCI conflict checks over metadata."),
    ("Admin & data health", " — manage users, monitor PM freshness, trigger syncs."),
]
runs = []
for b, rest in ops:
    runs.append([("▸  ", 13, CYAN, True), (b, 12.5, INK, True), (rest, 12.5, MUTED, False)])
text(s, Inches(0.6), Inches(2.05), Inches(7.0), Inches(4.2), runs, line_spacing=1.3, space_after=14)
tiles = [("🛰️", "Operational sites\nby tech & vendor"), ("🔌", "One-click external\nOSS / NMS launch"),
         ("🚗", "Drive-test trace\nvisualization"), ("❤️", "PM freshness\nmonitoring")]
gx, gy, cw, ch, gap = Inches(7.85), Inches(2.0), Inches(2.35), Inches(1.95), Inches(0.22)
for i, (ic, lbl) in enumerate(tiles):
    cx = gx + (i % 2) * (cw + gap)
    cy = gy + (i // 2) * (ch + gap)
    box(s, cx, cy, cw, ch, fill=CARD, line_color=LINE)
    text(s, cx, cy + Inches(0.28), cw, Inches(0.6), [[(ic, 30, CYAN, True)]], align=PP_ALIGN.CENTER)
    text(s, cx, cy + Inches(1.1), cw, Inches(0.7),
         [[(ln, 10.5, MUTED, False)] for ln in lbl.split("\n")], align=PP_ALIGN.CENTER, line_spacing=1.05)


# ============================================================
# SLIDE 11 — SECURITY
# ============================================================
s = add_slide(NAVY_DK)
eyebrow(s, "Security & Governance")
title(s, "Controlled access, full accountability")
sec = os.path.join(ASSETS, "security.png")
if os.path.exists(sec):
    img_box(s, sec, Inches(0.6), Inches(2.0), Inches(5.5), Inches(4.3))
secpts = [
    ("Role-based access", " — user → operator → NOC/systems → administrator."),
    ("Server-side sessions", " — httponly token cookie, expiry & active-user checks."),
    ("Password policy", " — 60-day rotation enforced via before-request gate."),
    ("CSRF + input sanitization", " on every state-changing request."),
    ("Activity logging", " & optional license activation for production."),
    ("Registration disabled", " — provisioned accounts only."),
]
runs = []
for b, rest in secpts:
    runs.append([("▸  ", 13, CYAN, True), (b, 12.5, INK, True), (rest, 12.5, MUTED, False)])
text(s, Inches(6.5), Inches(2.05), Inches(6.2), Inches(4.0), runs, line_spacing=1.25, space_after=12)
cx = Inches(6.5)
for lab in ["CSP headers", "SQLite WAL", "sync_log audit", "License gate"]:
    cx = chip(s, cx, Inches(5.95), lab)


# ============================================================
# SLIDE 12 — TECH STACK
# ============================================================
s = add_slide()
eyebrow(s, "Under the Hood")
title(s, "Technology stack")
tech = [
    ("🐍", "Flask", "Modular monolith; one blueprint per feature area."),
    ("🗄️", "SQLite", "WAL mode, busy-timeout, attach-based reads."),
    ("⏰", "APScheduler", "Watcher, hourly & daily ingest jobs."),
    ("🔐", "Sessions", "DB-backed tokens, salted hashing, roles."),
    ("📂", "SFTP / Paramiko", "Vendor PM, metadata & neighbor feeds."),
    ("🔄", "ncm_core", "Shared XML / Excel converters & comparators."),
    ("🐳", "Docker", "Web + scheduler services, Gunicorn, healthchecks."),
    ("🤖", "Optional AI", "OpenAI-backed parameter summaries (fallback safe)."),
]
cw, ch, gx, gy, gap = Inches(2.95), Inches(1.7), Inches(0.7), Inches(2.0), Inches(0.22)
for i, (ic, h, b) in enumerate(tech):
    cx = gx + (i % 4) * (cw + gap)
    cy = gy + (i // 4) * (ch + gap)
    card(s, cx, cy, cw, ch, ic, h, b, body_size=10)
text(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.6),
     [[("Deployed as two containers — ", 12, MUTED, False),
       ("primenet-web", 12, CYAN, True), (" (Gunicorn) and ", 12, MUTED, False),
       ("primenet-scheduler", 12, CYAN, True), (" — sharing one data volume.", 12, MUTED, False)]],
     align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13 — FEATURE CATALOG
# ============================================================
s = add_slide(NAVY_DK)
eyebrow(s, "Feature Catalog")
title(s, "19 blueprints, one coherent platform")
feats = [
    ("Performance", "KPI analytics engine"), ("Network Map", "Sites · wedges · KPIs"),
    ("Neighbor Analysis", "Relation lines"), ("Cell Heatmap", "Spatial KPI heat"),
    ("Conflict Map", "PCI risk scoring"), ("Femto PM", "Small-cell trends"),
    ("Reports", "Generate & archive"), ("Parameter Dict.", "MO lookup + search"),
    ("XML Parser", "XML → Excel"), ("Excel Generator", "Excel → XML"),
    ("NE Comparison", "Diff NE XMLs"), ("Config History", "Versions & diff"),
    ("Task Scheduler", "Config task flow"), ("RAN Features", "Feature reference"),
    ("Network Mgmt", "Sites & cells"), ("Drive-Test Viewer", "GPX / NMFS"),
    ("User Profile", "Prefs & activity"), ("Admin Panel", "Users & freshness"),
    ("Sync", "Pipeline triggers"), ("Auth", "Login · dashboard"),
]
cw, ch, gx, gy, gxgap, gygap = Inches(2.35), Inches(0.92), Inches(0.55), Inches(1.95), Inches(0.18), Inches(0.16)
for i, (h, b) in enumerate(feats):
    cx = gx + (i % 5) * (cw + gxgap)
    cy = gy + (i // 5) * (ch + gygap)
    c = box(s, cx, cy, cw, ch, fill=CARD, line_color=LINE)
    text(s, cx + Inches(0.14), cy + Inches(0.1), cw - Inches(0.24), ch - Inches(0.18),
         [[(h, 11.5, INK, True)], [(b, 9.5, MUTED, False)]], line_spacing=1.05)


# ============================================================
# SLIDE 14 — MATURITY / ROADMAP
# ============================================================
s = add_slide()
eyebrow(s, "Maturity — Honest Snapshot")
title(s, "Shipping today, building tomorrow")
cols = [
    (TEAL, "● PRODUCTION TODAY", "Live in operations",
     "KPI analytics · network map · neighbor analysis · reports · conflict map · Femto PM · parameter dictionary · XML/Excel tooling · NE compare · config tasks & history · dashboard · external tools · drive test · network mgmt · tasks."),
    (ROADBLUE, "◐ IN DEVELOPMENT", "Coming next",
     "Cell heatmap · SON analytics · network health dashboard."),
    (MUTED, "○ PLANNED", "On the roadmap",
     "RAN features catalog · OpenAPI reference · per-DB data dictionary · deployment playbook & incident runbooks."),
]
cw, ch, gap = Inches(3.95), Inches(4.0), Inches(0.22)
x = Inches(0.6)
for ac, tag, head, body in cols:
    c = box(s, x, Inches(2.0), cw, ch, fill=CARD, line_color=ac, line_w=1.5)
    text(s, x + Inches(0.22), Inches(2.2), cw - Inches(0.44), ch - Inches(0.4),
         [[(tag, 11, ac, True)],
          [("", 4, ac, False)],
          [(head, 15, INK, True)],
          [("", 4, ac, False)],
          [(body, 11.5, MUTED, False)]], line_spacing=1.25, space_after=4)
    x += cw + gap


# ============================================================
# SLIDE 15 — OUTCOMES
# ============================================================
s = add_slide(NAVY_DK)
eyebrow(s, "Outcomes & Value")
title(s, "Measurable wins — without displacing vendor investments")
out = [
    ("⚡", "Time to insight", "KPI trends & maps minutes after data lands — fewer manual exports."),
    ("🤝", "Cross-vendor consistency", "One platform, one set of KPI definitions for Nokia & Huawei."),
    ("👁️", "Operational visibility", "Live footprint by RAT & vendor; admins watch data freshness."),
    ("🛡️", "Configuration quality", "Structured compare / convert reduces undetected config drift."),
    ("🧹", "Less tool sprawl", "A single entry point plus curated launchers for daily systems."),
    ("📜", "Governance", "Roles, session security, password policy, activity logging."),
]
cw, ch, gx, gy, gap = Inches(5.95), Inches(1.55), Inches(0.7), Inches(2.0), Inches(0.25)
for i, (ic, h, b) in enumerate(out):
    cx = gx + (i % 2) * (cw + gap)
    cy = gy + (i // 2) * (ch + gap)
    card(s, cx, cy, cw, ch, ic, h, b, body_size=11)


# ============================================================
# SLIDE 16 — CLOSING
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
solid(bg, NAVY_DK); bg.shadow.inherit = False
if os.path.exists(hero):
    img_box(s, hero, 0, 0, EMU_W, EMU_H)
veil = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
solid(veil, NAVY_DK, alpha=78); veil.shadow.inherit = False
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Pt(4))
solid(bar, CYAN); bar.shadow.inherit = False
eyebrow(s, "Next Steps", l=Inches(5.7), t=Inches(1.2))
text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.0),
     [[("One front door for RAN ops", 40, CYAN, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(1.6), Inches(2.85), Inches(10.1), Inches(1.2),
     [[("PrimeNet turns fragmented radio operations into a single, governed workspace for", 15, INK, True)],
      [("performance intelligence, configuration support, and daily tool access.", 15, INK, True)]],
     align=PP_ALIGN.CENTER, line_spacing=1.3)
nxt = [("1 · Endorse", "Adopt PrimeNet as the standard internal entry point."),
       ("2 · Align", "Set data-freshness SLAs — hourly PM, daily deep sync."),
       ("3 · Prioritize", "Network health, SON analytics, cell heatmap.")]
cw, gap = Inches(3.85), Inches(0.25)
x = int((EMU_W - (cw * 3 + gap * 2)) / 2)
for h, b in nxt:
    c = box(s, x, Inches(4.5), cw, Inches(1.5), fill=RGBColor(0x10, 0x1B, 0x30), line_color=LINE)
    text(s, x + Inches(0.2), Inches(4.7), cw - Inches(0.4), Inches(1.1),
         [[(h, 15, INK, True)], [(b, 11, MUTED, False)]], align=PP_ALIGN.CENTER, line_spacing=1.15, space_after=6)
    x += cw + gap
footer(s, "PrimeNet · Radio Network Performance & Configuration Platform · June 2026")


out_path = os.path.join(HERE, "PrimeNet_Presentation.pptx")
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides._sldIdLst))
